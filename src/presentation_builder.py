# ===================================================================================
# FILE 3: src/presentation_builder.py
# ===================================================================================

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
from scatter_plot import ScatterPlotGenerator

class PresentationBuilder:
    """Builds PowerPoint presentation with Lapin branding"""

    def __init__(self, processor, chart_gen, assets_path='assets'):
        self.processor = processor
        self.chart_gen = chart_gen
        self.scatter_gen = ScatterPlotGenerator(processor)
        self.assets_path = assets_path
        self.prs = Presentation()
        
        # FIXED DIMENSIONS: Widescreen 16:9
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        
        # Logo paths (corrected file names)
        self.logo_lockup_white = os.path.join(assets_path, 'logos', 'Lapin_LogoLockup_RGB_White.png')
        self.logo_lockup_gray = os.path.join(assets_path, 'logos', 'Lapin_LogoLockup_RGB_DarkGray.png')
        self.logo_symbol_white = os.path.join(assets_path, 'logos', 'Lapin_Symbol_RGB_White.png')
        self.logo_symbol_navy = os.path.join(assets_path, 'logos', 'Lapin_Symbol_RGB_Navy.png')

    def _get_ordered_categories(self):
        if hasattr(self.processor, "get_ordered_categories"):
            return self.processor.get_ordered_categories()
        return self.processor.get_all_categories()

    def _get_ordered_roles(self):
        if hasattr(self.processor, "get_ordered_roles"):
            return self.processor.get_ordered_roles()
        return list(self.processor.roles)
    
    def _add_logo(self, slide, logo_path, left, top, height):
        """Add logo to slide if file exists"""
        if os.path.exists(logo_path):
            slide.shapes.add_picture(logo_path, left, top, height=height)
            return True
        return False
    
    def _set_font(self, text_frame, font_name='Ginto Normal Medium', font_size=16, color_rgb=(26, 26, 66)):
        """Apply Ginto font to text"""
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = font_name
                run.font.size = Pt(font_size)
                run.font.color.rgb = RGBColor(*color_rgb)
    
    def add_title_slide(self, title="Lapin Trust Index", subtitle="Sample Report"):
        """Title slide with rowing background"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Add rowing background image
        import os
        rowing_img = os.path.join(self.assets_path, 'images', 'rowing_team.jpg')
        if os.path.exists(rowing_img):
            slide.shapes.add_picture(rowing_img, Inches(0), Inches(0), 
                            width=self.prs.slide_width, 
                            height=self.prs.slide_height)
        else:
            # Fallback to navy background
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(26, 26, 66)
        
        # White logo lockup (top left)
        self._add_logo(slide, self.logo_lockup_white, Inches(0.5), Inches(0.5), Inches(0.7))
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.), Inches(11), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_p = title_frame.paragraphs[0]
        title_p.font.name = 'Ginto Normal Medium'
        title_p.font.size = Pt(54)
        title_p.font.bold = True
        title_p.alignment = PP_ALIGN.LEFT
        title_p.font.color.rgb = RGBColor(255, 255, 255)
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(8), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_p = subtitle_frame.paragraphs[0]
        subtitle_p.font.name = 'Ginto Normal Medium'
        subtitle_p.font.size = Pt(24)
        subtitle_p.font.color.rgb = RGBColor(255, 255, 255)
    
    def add_methodology_slide(self):
        """Methodology slide with white symbol logo"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Navy background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(26, 26, 66)
        
        # White symbol only (NOT lockup)
        self._add_logo(slide, self.logo_symbol_white, Inches(12.3), Inches(6.5), Inches(0.6))
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "Methodology"
        title_p = title_frame.paragraphs[0]
        title_p.font.name = 'Ginto Normal Medium'
        title_p.font.size = Pt(48)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(255, 255, 255)
        
        # Content with proper wrapping
        categories = self._get_ordered_categories()
        category_bullets = '\n    • '.join(categories)

        content = f"""Questions focused on five key dimensions that comprise the Lapin Performance Index:
    • {category_bullets}

• Scores are delineated by rater group: managers and directors/VPs
• Likert scale ratings are shown grouped into "agree" responses—those with a 4 or 5 rating.
• We analyzed patterns, interpreted results, and identified key insights to help move the team forward."""

        content_box = slide.shapes.add_textbox(Inches(0.72), Inches(1.7), Inches(7), Inches(4.5))

        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        content_frame.text = content
        
        for paragraph in content_frame.paragraphs:
            paragraph.font.name = 'Ginto Normal Medium'
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
        
        # Footer bottom left + page number
        footer_box = slide.shapes.add_textbox(Inches(0.7), Inches(6.9), Inches(6), Inches(0.3))
        footer_frame = footer_box.text_frame
        footer_frame.text = "2     Timeless Wisdom. Timely Solutions."
        footer_p = footer_frame.paragraphs[0]
        footer_p.font.name = 'Ginto Normal Medium'
        footer_p.font.size = Pt(10)
        footer_p.font.color.rgb = RGBColor(255, 255, 255)

        # Add decorative graphics (right side)
        # Navy lines (left graphic)
        navy_lines = os.path.join(self.assets_path, 'images', 'navy_lines.png')
        if os.path.exists(navy_lines):
            slide.shapes.add_picture(navy_lines, Inches(11.7), Inches(4.5), height=Inches(1.5))

        # White circle (right graphic)
        circle_graphic = os.path.join(self.assets_path, 'images', 'white_circle.png')
        if os.path.exists(circle_graphic):
            slide.shapes.add_picture(circle_graphic, Inches(10.2), Inches(4.5), height=Inches(1.5))

    def add_section_divider(self, title):
        """Category divider slide with gray logo lockup"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Light gray background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(208, 211, 214)  # Light Gray
        
        # FIXED: Add gray logo lockup (top left)
        self._add_logo(slide, self.logo_lockup_gray, Inches(0.5), Inches(0.5), Inches(1))
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(6), Inches(2.5))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_frame.text = title
        title_p = title_frame.paragraphs[0]
        title_p.font.name = 'Ginto Normal Medium'
        title_p.font.size = Pt(54)
        title_p.font.color.rgb = RGBColor(139, 145, 151)  # Dark Gray
        title_p.alignment = PP_ALIGN.LEFT
    
        small_white_circle = os.path.join(self.assets_path, 'images', 'small_white_circle.png')
        if os.path.exists(small_white_circle):
        # Add picture
            circle_size = self.prs.slide_height  # Full height (7.5")
            left_position = self.prs.slide_width - (circle_size / 2)  # Center on right edge
    
            pic = slide.shapes.add_picture(small_white_circle, left_position, Inches(0), 
                                   height=circle_size)
        # Rotate 270 degrees clockwise
            pic.rotation = 270

    def add_insights_slide(self, slide_number=None):
        """Ivory insights and findings slide with navy symbol logo"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Ivory background (slides 34 and 97 only)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(247, 245, 242)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "Insights and Findings"
        title_p = title_frame.paragraphs[0]
        title_p.font.name = 'Ginto Normal Medium'
        title_p.font.size = Pt(36)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(26, 26, 66)  # Navy

        # Placeholder content
        content_box = slide.shapes.add_textbox(Inches(0.7), Inches(2), Inches(5), Inches(1))
        content_frame = content_box.text_frame
        content_frame.text = "Analysis..."
        content_p = content_frame.paragraphs[0]
        content_p.font.name = 'Ginto Normal Medium'
        content_p.font.size = Pt(18)
        content_p.font.bold = True
        content_p.font.color.rgb = RGBColor(26, 26, 66)

        # Second placeholder
        content_box2 = slide.shapes.add_textbox(Inches(6.5), Inches(2), Inches(5), Inches(1))
        content_frame2 = content_box2.text_frame
        content_frame2.text = "Analysis..."
        content_p2 = content_frame2.paragraphs[0]
        content_p2.font.name = 'Ginto Normal Medium'
        content_p2.font.size = Pt(18)
        content_p2.font.bold = True
        content_p2.font.color.rgb = RGBColor(26, 26, 66)

        # Footer + page number
        page_num = slide_number if slide_number else ""
        footer_box = slide.shapes.add_textbox(Inches(0.7), Inches(6.9), Inches(6), Inches(0.3))
        footer_frame = footer_box.text_frame
        footer_frame.text = f"{page_num}     Timeless Wisdom. Timely Solutions."
        footer_p = footer_frame.paragraphs[0]
        footer_p.font.name = 'Ginto Normal Medium'
        footer_p.font.size = Pt(10)
        footer_p.font.color.rgb = RGBColor(26, 26, 66)

        # Navy symbol logo (bottom right)
        self._add_logo(slide, self.logo_symbol_navy, Inches(12.3), Inches(6.5), Inches(0.6))

    def add_chart_slide(self, title, chart_bytes, slide_number=None, chart_left=None, chart_top=None, chart_width=None, chart_height=Inches(5.5)):
        """Chart slide with navy symbol logo (bottom RIGHT)"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Ivory background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(247, 245, 242)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12), Inches(1.2))
        title_frame = title_box.text_frame
        title_frame.word_wrap = False
        title_frame.text = title
        title_p = title_frame.paragraphs[0]
        title_p.font.name = 'Ginto Normal Medium'
        title_p.font.size = Pt(30)
        title_p.font.color.rgb = RGBColor(26, 26, 66)  # Navy

        # Chart image — Option 2 spec: 12" wide × 5.5" tall at (0.75", 1.51")
        # chart_height=None lets python-pptx auto-scale height from native image aspect ratio
        left = chart_left if chart_left is not None else Inches(0.75)
        top = chart_top if chart_top is not None else Inches(1.51)
        width = chart_width if chart_width is not None else Inches(12)
        if chart_height is not None:
            slide.shapes.add_picture(chart_bytes, left, top, width=width, height=chart_height)
        else:
            slide.shapes.add_picture(chart_bytes, left, top, width=width)
        
        # Footer + page number (BOTTOM LEFT)
        if slide_number:
            footer_box = slide.shapes.add_textbox(Inches(0.7), Inches(6.9), Inches(6), Inches(0.3))
            footer_frame = footer_box.text_frame
            footer_frame.text = f"{slide_number}     Timeless Wisdom. Timely Solutions."
            footer_p = footer_frame.paragraphs[0]
            footer_p.font.name = 'Ginto Normal Medium'
            footer_p.font.size = Pt(10)
            footer_p.font.color.rgb = RGBColor(26, 26, 66)
        
        # Navy symbol logo (BOTTOM RIGHT)
        self._add_logo(slide, self.logo_symbol_navy, Inches(12.3), Inches(6.5), Inches(0.6))
    
    def build_full_report(self):
        """Generate complete 97-slide report following Lapin's exact structure"""
        print("\n" + "="*60)
        print("BUILDING POWERPOINT PRESENTATION")
        print("="*60)
        
        slide_count = 1
        
        # ===================================================================
        # SECTION 1: INTRODUCTION (Slides 1-2)
        # ===================================================================
        print(f"\n[{slide_count}] Title slide")
        self.add_title_slide()
        slide_count += 1
        
        print(f"[{slide_count}] Methodology")
        self.add_methodology_slide()
        slide_count += 1
        
        # ===================================================================
        # SECTION 2: CATEGORY OVERVIEW (Slides 3-32)
        # Pattern: Divider + Multi-line + Overall dist + 3 role dists = 6 slides per category
        # ===================================================================
        categories = self._get_ordered_categories()
        
        print("\n" + "="*60)
        print("SECTION 2: CATEGORY OVERVIEW")
        print("="*60)
        
        for category in categories:
            questions = self.processor.get_category_questions(category)
            if len(questions) == 0:
                continue
            
            print(f"\n--- {category} (6 slides) ---")
            
            # 1. Section divider
            print(f"[{slide_count}] {category} Scores (divider)")
            self.add_section_divider(f"{category} Scores")
            slide_count += 1
            
            # 2. Multi-line chart (all roles)
            print(f"[{slide_count}] {category} Scores: Percent Agree (multi-line)")
            chart = self.chart_gen.create_multi_line(questions, category)
            self.add_chart_slide(
                f'{category} Scores: Percent "Agree" (4 or 5 rating)', 
                chart, slide_count
            )
            slide_count += 1
            
            # 3. Overall distribution
            print(f"[{slide_count}] Overall {category} Distribution")
            chart = self.chart_gen.create_stacked_bar(questions, role=None)
            self.add_chart_slide(f'Overall {category} Distribution', chart, slide_count)
            slide_count += 1
            
            # 4-6. Distribution for each role (THIS WAS MISSING!)
            for role in self._get_ordered_roles():
                print(f"[{slide_count}] {role}: {category} Distribution")
                chart = self.chart_gen.create_stacked_bar(questions, role)
                self.add_chart_slide(f'{role}: {category} Distribution', chart, slide_count)
                slide_count += 1
        
        # ===================================================================
        # SECTION 3: FINDINGS INTRO (Slides 33-34)
        # ===================================================================
        print("\n" + "="*60)
        print("SECTION 3: FINDINGS INTRO")
        print("="*60)
        
        print(f"[{slide_count}] Findings (divider)")
        self.add_section_divider("Findings")
        slide_count += 1
        
        print(f"[{slide_count}] Insights and Findings (placeholder)")
        # TODO: This will be a custom slide with specific content
        self.add_insights_slide(slide_count)
        slide_count += 1

        
        # ===================================================================
        # SECTION 4-6: ROLE DEEP DIVES (19 slides each × 3 roles = 57 slides)
        # ===================================================================
        print("\n" + "="*60)
        print("SECTIONS 4-6: ROLE DEEP DIVES")
        print("="*60)
        
        for role in self._get_ordered_roles():
            print(f"\n{'='*60}")
            print(f"ROLE: {role} (19 slides)")
            print('='*60)
            
            # 1. "{Role} Findings" divider
            print(f"[{slide_count}] {role} Findings (divider)")
            self.add_section_divider(f"{role} Findings")
            slide_count += 1
            
            # 2. "{Role}: Averages" divider
            print(f"[{slide_count}] {role}: Averages (divider)")
            self.add_section_divider(f"{role}: Averages")
            slide_count += 1
            
            # 3-7. Average bars for each category (5 slides)
            print(f"\n  Average bars by category:")
            for category in categories:
                questions = self.processor.get_category_questions(category)
                if len(questions) == 0:
                    continue
                
                print(f"  [{slide_count}] {role}: {category} Average")
                chart = self.chart_gen.create_average_bars(questions, role)
                self.add_chart_slide(f'{role}: {category} Average', chart, slide_count)
                slide_count += 1
            
            # 8. "{Role} Agree Ratings" divider
            print(f"\n[{slide_count}] {role} Agree Ratings (4 or 5 responses) (divider)")
            self.add_section_divider(f"{role} Agree Ratings (4 or 5 responses)")
            slide_count += 1
            
            # 9-13. Single-line % agree for each category (5 slides)
            print(f"\n  Agree ratings by category:")
            for category in categories:
                questions = self.processor.get_category_questions(category)
                if len(questions) == 0:
                    continue
                
                print(f"  [{slide_count}] {role}: {category} Percent Agree")
                chart = self.chart_gen.create_single_line(questions, role)
                self.add_chart_slide(
                    f'{role}: {category} Percent "Agree" (4 or 5 rating)', 
                    chart, slide_count
                )
                slide_count += 1
            
            # 14. "{Role} Distribution per Category" divider
            print(f"\n[{slide_count}] {role} Distribution per Category (divider)")
            self.add_section_divider(f"{role} Distribution per Category")
            slide_count += 1
            
            # 15-19. Stacked distributions for each category (5 slides) - DUPLICATES
            print(f"\n  Distributions by category (duplicates from Section 2):")
            for category in categories:
                questions = self.processor.get_category_questions(category)
                if len(questions) == 0:
                    continue
                
                print(f"  [{slide_count}] {role}: {category} Distribution (DUPLICATE)")
                chart = self.chart_gen.create_stacked_bar(questions, role)
                self.add_chart_slide(f'{role}: {category} Distribution', chart, slide_count)
                slide_count += 1
        
        # ===================================================================
        # SECTION 7: HIGH PERFORMING TEAMS (Slides 92-95)
        # ===================================================================
        print("\n" + "="*60)
        print("SECTION 7: HIGH PERFORMING TEAMS")
        print("="*60)

        print(f"[{slide_count}] High Performing Teams (divider)")
        self.add_section_divider("High Performing Teams")
        slide_count += 1

        for role in self._get_ordered_roles():
            print(f"[{slide_count}] High Performing Teams – {role}")
            chart = self.scatter_gen.create(role)
            self.add_chart_slide(f'High Performing Teams – {role}', chart, slide_count,
                                 chart_left=Inches(0.667), chart_top=Inches(1.5), chart_width=Inches(12),
                                 chart_height=None)
            slide_count += 1
        
        # ===================================================================
        # SECTION 8: CLOSING (Slides 96-97)
        # ===================================================================
        print("\n" + "="*60)
        print("SECTION 8: CLOSING")
        print("="*60)
        
        print(f"[{slide_count}] Findings (divider - duplicate)")
        self.add_section_divider("Findings")
        slide_count += 1
        
        print(f"[{slide_count}] Insights and Findings (placeholder)")
        self.add_insights_slide(slide_count)
        slide_count += 1
        
        print("\n" + "="*60)
        print(f"✓ PRESENTATION COMPLETE: {slide_count - 1} slides generated")
        print("="*60)
        
        return self.prs
